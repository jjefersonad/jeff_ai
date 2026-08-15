"""Testes do PptxWriter (infrastructure/pptx_writer) e da tool create_pptx_presentation.

Cobrem os critérios de aceitação da task `custom-office-doc-tools-task-pptx-1`:
- REQ-001: gera `.pptx` válido só com python-pptx, sem binários externos.
- REQ-002: suporta slide de título, slides de conteúdo com bullets, inserção
  de imagem e tabela simples.
- REQ-003: salva em `outputs/documents/pptx/<ts>.pptx`, retorna
  `{path, url, metadata}` e a URL começa com `/api/files/pptx/`.
- REQ-005: entrada inválida retorna erro descritivo e não deixa arquivo parcial.
- Tool como adapter fino + wiring em `composition.dependencies`.
"""
from __future__ import annotations

import struct
import zlib
from pathlib import Path

import pytest
from pptx import Presentation

import src.composition.dependencies as dep
import src.tools.create_pptx_presentation_tool as pptx_tool
from src.domain.documents import (
    BulletSlide,
    ImageRef,
    ImageSlide,
    PptxSpec,
    Table,
    TableSlide,
    TitleSlide,
)
from src.domain.shared.errors import DomainError
from src.infrastructure.documents.pptx_writer import PptxWriter
from src.models.pptx_document import PptxDocumentInput, PptxSlideInput

# --- helpers ---------------------------------------------------------------


def _make_png(path: Path, width: int = 2, height: int = 2) -> str:
    """Grava um PNG mínimo válido (2x2) no caminho e retorna o caminho como str.

    Evita dependência de Pillow nos testes — gera bytes PNG manualmente
    (magic + IHDR + IDAT com pixels zero + IEND).
    """

    def _chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 0, 0, 0, 0)
    raw = b"".join(b"\x00" + b"\x00\x00\x00\x00" * width for _ in range(height))
    idat = zlib.compress(raw)
    png_bytes = sig + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", idat) + _chunk(b"IEND", b"")
    path.write_bytes(png_bytes)
    return str(path)


# --- PptxWriter (infra) ----------------------------------------------------


@pytest.fixture
def writer(tmp_path: Path) -> PptxWriter:
    """PptxWriter apontando para um diretório temporário (sem poluir outputs/)."""
    return PptxWriter(output_dir=tmp_path, url_prefix="/api/files")


async def test_writer_creates_valid_pptx_with_title_and_bullets(writer, tmp_path):
    """REQ-001/REQ-002/REQ-003: writer gera .pptx válido, capa + bullets."""
    spec = PptxSpec(
        slides=(
            TitleSlide(title="Roadmap 2026", subtitle="Time de IA"),
            BulletSlide(
                title="Objetivos",
                bullets=("Reduzir latência", "Aumentar cobertura", "Manter qualidade"),
            ),
        )
    )

    result = await writer.write(spec)

    # REQ-003: contrato {path, url, metadata}, URL pública e arquivo no disco.
    assert result.path
    assert result.url.startswith("/api/files/pptx/")
    assert result.url.endswith(".pptx")
    assert result.metadata == {"kind": "pptx", "slide_count": 2}
    path = Path(result.path)
    assert path.is_file()
    assert path.parent == tmp_path

    # Reabre com python-pptx e valida a estrutura.
    presentation = Presentation(str(path))
    assert len(presentation.slides) == 2

    slide0 = presentation.slides[0]
    assert slide0.shapes.title.text == "Roadmap 2026"

    slide1 = presentation.slides[1]
    assert slide1.shapes.title.text == "Objetivos"
    # Encontra o placeholder de body e valida os bullets.
    body_texts = []
    for shape in slide1.shapes:
        if shape.has_text_frame and shape != slide1.shapes.title:
            body_texts.extend(p.text for p in shape.text_frame.paragraphs)
    assert body_texts == ["Reduzir latência", "Aumentar cobertura", "Manter qualidade"]


async def test_writer_embeds_image_in_slide(writer, tmp_path):
    """REQ-002: slide com imagem existente em disco é embutido no .pptx."""
    image_path = tmp_path / "team.png"
    _make_png(image_path, width=4, height=4)

    spec = PptxSpec(
        slides=(
            ImageSlide(image=ImageRef(path=str(image_path), width_inches=3.0), title="Time"),
        )
    )

    result = await writer.write(spec)
    presentation = Presentation(result.path)

    assert len(presentation.slides) == 1
    slide = presentation.slides[0]
    # Há uma imagem no slide (Picture shape).
    pictures = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    assert len(pictures) == 1


async def test_writer_renders_table_slide_with_header_bold(writer, tmp_path):
    """REQ-002: slide com tabela simples; primeira linha em negrito quando header=True."""
    spec = PptxSpec(
        slides=(
            TableSlide(
                table=Table(
                    rows=(("Mês", "Receita"), ("Jan", "1000"), ("Fev", "1200")),
                    header=True,
                ),
                title="Vendas",
            ),
        )
    )

    result = await writer.write(spec)
    presentation = Presentation(result.path)
    assert len(presentation.slides) == 1

    # Encontra a tabela no slide.
    tables = [s.table for s in presentation.slides[0].shapes if s.has_table]
    assert len(tables) == 1
    table = tables[0]
    assert len(table.rows) == 3
    assert len(table.columns) == 2
    assert table.cell(0, 0).text == "Mês"
    # Cabeçalho em negrito.
    header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert header_run.font.bold is True
    # Dados: não-negrito.
    body_run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
    assert body_run.font.bold is None or body_run.font.bold is False


async def test_writer_missing_image_raises_without_partial_file(writer, tmp_path):
    """REQ-005: imagem inexistente levanta erro claro e nada fica no diretório."""
    spec = PptxSpec(
        slides=(
            ImageSlide(image=ImageRef(path="/caminho/que/nao/existe.png")),
        )
    )

    with pytest.raises(RuntimeError, match="Imagem não encontrada"):
        await writer.write(spec)

    # Sem arquivo parcial.
    assert list(tmp_path.iterdir()) == []


async def test_writer_rejects_non_pptx_spec():
    """O writer só aceita PptxSpec — typecheck em tempo de execução."""
    writer = PptxWriter()
    with pytest.raises(TypeError, match="PptxSpec"):
        await writer.write("not a spec")  # type: ignore[arg-type]


# --- Domínio: validações que protegem o writer ----------------------------


def test_pptx_spec_rejects_empty_slides():
    with pytest.raises(DomainError, match="ao menos um slide"):
        PptxSpec(slides=())


def test_title_slide_rejects_empty_title():
    with pytest.raises(DomainError, match="title"):
        TitleSlide(title="")


def test_title_slide_rejects_empty_subtitle_when_informed():
    with pytest.raises(DomainError, match="subtitle"):
        TitleSlide(title="ok", subtitle="   ")


def test_bullet_slide_rejects_empty_bullets():
    with pytest.raises(DomainError, match="bullets"):
        BulletSlide(title="ok", bullets=())


def test_bullet_slide_rejects_empty_bullet_string():
    with pytest.raises(DomainError, match="bullets"):
        BulletSlide(title="ok", bullets=("a", ""))


# --- Tool create_pptx_presentation (HTML pipeline) -------------------------


def test_slides_to_html_title_and_bullets():
    html = pptx_tool.slides_to_html(
        [
            PptxSlideInput(type="title", title="T", subtitle="S"),
            PptxSlideInput(type="bullets", title="B", bullets=["a", "b"]),
        ]
    )
    assert 'class="slide"' in html
    assert "<h1>T</h1>" in html
    assert "<p>S</p>" in html
    assert "<h2>B</h2>" in html
    assert "<li>a</li>" in html


def test_slides_to_html_table_and_image_title():
    html = pptx_tool.slides_to_html(
        [
            PptxSlideInput(type="table", title="Tab", rows=[["a", "b"]], header=False),
            PptxSlideInput(type="image", title="Img", path="/tmp/x.png"),
        ]
    )
    assert "<h2>Tab</h2>" in html
    assert "<td>a</td>" in html
    assert "<h2>Img</h2>" in html


def test_slides_to_html_drops_incomplete_slides():
    """Slides incompletos são descartados; se nada sobrar → DomainError."""
    from src.domain.shared.errors import DomainError

    with pytest.raises(DomainError, match="slide"):
        pptx_tool.slides_to_html(
            [
                PptxSlideInput(type="title", title=""),
                PptxSlideInput(type="bullets", title="T"),
                PptxSlideInput(type="image", path="/tmp/x.png"),  # sem título
                PptxSlideInput(type="table", title="T"),  # sem rows
                PptxSlideInput(type="unknown", title="Z"),
            ]
        )


def test_slides_to_html_skips_invalid_keeps_valid():
    html = pptx_tool.slides_to_html(
        [
            PptxSlideInput(type="title", title=""),
            PptxSlideInput(type="title", title="Capa"),
            PptxSlideInput(type="bullets", title="T", bullets=["x"]),
            PptxSlideInput(type="unknown", title="Z"),
        ]
    )
    assert html.count('class="slide"') == 2


async def test_tool_returns_path_url_metadata(monkeypatch, tmp_path):
    """A tool usa o pipeline HTML e devolve o contrato esperado."""
    from unittest.mock import AsyncMock

    from src.infrastructure.documents.html_pptx_converter import HtmlPptxConverter

    monkeypatch.setattr(pptx_tool, "require_user_docs_dir", AsyncMock(return_value=tmp_path))
    monkeypatch.setattr(pptx_tool, "_document_url_prefix", lambda: "/api/files")
    monkeypatch.setattr(pptx_tool, "record_ownership", AsyncMock())

    render = pptx_tool._build_pptx_render(tmp_path)
    assert isinstance(render._converters["pptx"], HtmlPptxConverter)

    result = await pptx_tool.create_pptx_presentation.coroutine(
        PptxDocumentInput(
            slides=[
                PptxSlideInput(type="title", title="Capa"),
                PptxSlideInput(type="bullets", title="T", bullets=["a", "b"]),
            ],
        )
    )

    assert set(result) == {"path", "url", "metadata"}
    assert "/api/files/pptx/" in result["url"]
    assert result["url"].endswith(".pptx")
    assert result["metadata"]["kind"] == "pptx"
    assert Path(result["path"]).is_file()


async def test_tool_returns_error_when_all_slides_invalid(monkeypatch, tmp_path):
    """REQ-005: nenhum slide válido → `error` descritivo, sem arquivo parcial."""
    from unittest.mock import AsyncMock

    monkeypatch.setattr(pptx_tool, "require_user_docs_dir", AsyncMock(return_value=tmp_path))
    monkeypatch.setattr(pptx_tool, "_document_url_prefix", lambda: "/api/files")
    monkeypatch.setattr(pptx_tool, "record_ownership", AsyncMock())

    result = await pptx_tool.create_pptx_presentation.coroutine(
        PptxDocumentInput(
            slides=[
                PptxSlideInput(type="title", title=""),
                PptxSlideInput(type="bullets", title="T"),
            ],
        )
    )

    assert "error" in result
    assert "slide" in result["error"].lower()
    assert list(tmp_path.iterdir()) == []
