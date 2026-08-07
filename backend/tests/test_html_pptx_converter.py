"""Testes do converter HTML→PPTX (html-document-tools-task-pptx-1)."""
from __future__ import annotations

from io import BytesIO

import pytest
from pptx import Presentation

from src.domain.shared.errors import DomainError
from src.infrastructure.documents.html_pptx_converter import HtmlPptxConverter


@pytest.mark.asyncio
async def test_section_slide_becomes_pptx_slide() -> None:
    html = (
        '<section class="slide">'
        "<h1>Título Principal</h1>"
        "<p>Introdução.</p>"
        "</section>"
        '<div class="slide">'
        "<h2>Segundo</h2>"
        "<ul><li>A</li><li>B</li></ul>"
        "</div>"
    )
    converter = HtmlPptxConverter()
    payload = await converter.convert(html=html, css=None, kind="pptx")

    assert payload[:2] == b"PK"
    prs = Presentation(BytesIO(payload))
    assert len(prs.slides) == 2
    texts = [
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]
    joined = "\n".join(texts)
    assert "Título Principal" in joined
    assert "Segundo" in joined


@pytest.mark.asyncio
async def test_html_without_slides_raises() -> None:
    converter = HtmlPptxConverter()
    with pytest.raises(DomainError, match="slide"):
        await converter.convert(html="<p>sem slides</p>", css=None, kind="pptx")
