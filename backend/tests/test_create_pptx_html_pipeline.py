"""create_pptx_presentation via HTML pipeline (html-document-tools-task-pptx-1)."""
from __future__ import annotations

from pathlib import Path
from unittest.mock import AsyncMock

import pytest
from pptx import Presentation

import src.tools.create_pptx_presentation_tool as pptx_tool
from src.models.html_document_input import HtmlDocumentInput
from src.models.pptx_document import PptxDocumentInput, PptxSlideInput


@pytest.fixture
def documents_root(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    root = tmp_path / "documents"
    monkeypatch.setattr(pptx_tool, "require_user_docs_dir", AsyncMock(return_value=root))
    monkeypatch.setattr(
        pptx_tool,
        "_document_url_prefix",
        lambda: "http://localhost:3000/api/files",
    )
    monkeypatch.setattr(pptx_tool, "record_ownership", AsyncMock())
    return root


async def test_html_slides_return_pptx_url(documents_root: Path) -> None:
    """Unit-1: section.slide com título → url /api/files/pptx/ + kind=pptx."""
    out = await pptx_tool.create_pptx_presentation.coroutine(
        HtmlDocumentInput(
            html=(
                '<section class="slide">'
                "<h1>Abertura</h1>"
                "<p>Bem-vindos</p>"
                "</section>"
            ),
            title="Deck",
        )
    )

    assert "error" not in out
    assert out["metadata"]["kind"] == "pptx"
    assert "/api/files/pptx/" in out["url"]
    path = Path(out["path"])
    assert path.is_file()
    assert path.read_bytes()[:2] == b"PK"

    prs = Presentation(str(path))
    assert len(prs.slides) >= 1
    texts = [shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame]
    assert any("Abertura" in t for t in texts)


async def test_no_slides_rejects_without_file(documents_root: Path) -> None:
    """Unit-2: HTML sem slides → error, sem pptx."""
    out = await pptx_tool.create_pptx_presentation.coroutine(
        HtmlDocumentInput(html="<p>Só texto, sem slides.</p>", title="Vazio")
    )

    assert "error" in out
    assert "path" not in out
    pptx_dir = documents_root
    assert not pptx_dir.exists() or list(pptx_dir.glob("*.pptx")) == []


async def test_legacy_slides_still_work(documents_root: Path) -> None:
    out = await pptx_tool.create_pptx_presentation.coroutine(
        PptxDocumentInput(
            slides=[PptxSlideInput(type="title", title="Capa", subtitle="Sub")]
        )
    )
    assert "error" not in out
    assert out["metadata"]["kind"] == "pptx"
    prs = Presentation(out["path"])
    assert len(prs.slides) >= 1
