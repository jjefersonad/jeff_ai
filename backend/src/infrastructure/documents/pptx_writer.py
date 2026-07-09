"""Adapter de criação de apresentações `.pptx` via python-pptx (implementa o port)."""
from __future__ import annotations

import asyncio
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from src.application.ports.document_writer import DocumentWriterPort
from src.domain.documents import (
    BulletSlide,
    DocumentResult,
    DocumentSpec,
    ImageRef,
    ImageSlide,
    PptxSpec,
    Table,
    TableSlide,
    TitleSlide,
)
from src.infrastructure.documents.output_target import DocumentOutput


class PptxWriter(DocumentWriterPort):
    """Gera um `.pptx` a partir de um `PptxSpec` usando apenas python-pptx.

    Não invoca soffice/pandoc/node. Persiste em `outputs/documents/pptx/` e monta
    a URL pública (`/api/files/pptx/...`). A apresentação é montada em memória e
    só é escrita no final (`save`), então falhas durante a montagem não deixam
    arquivo parcial.
    """

    def __init__(
        self,
        *,
        output_dir: Path | None = None,
        url_prefix: str = "/api/files",
    ) -> None:
        """Configura o destino e o prefixo de URL do writer."""
        self._output = DocumentOutput("pptx", output_dir=output_dir, url_prefix=url_prefix)

    async def write(self, spec: DocumentSpec) -> DocumentResult:
        """Gera o `.pptx` do `spec` (fora do event loop) e retorna o resultado."""
        if not isinstance(spec, PptxSpec):
            raise TypeError("PptxWriter só aceita PptxSpec.")
        return await asyncio.to_thread(self._write_sync, spec)

    def _write_sync(self, spec: PptxSpec) -> DocumentResult:
        presentation = Presentation()

        for slide in spec.slides:
            if isinstance(slide, TitleSlide):
                self._render_title_slide(presentation, slide)
            elif isinstance(slide, BulletSlide):
                self._render_bullet_slide(presentation, slide)
            elif isinstance(slide, ImageSlide):
                self._render_image_slide(presentation, slide)
            elif isinstance(slide, TableSlide):
                self._render_table_slide(presentation, slide)

        path, url = self._output.allocate(spec.extension)
        presentation.save(str(path))
        return DocumentResult(path=str(path), url=url, metadata=spec.metadata())

    @staticmethod
    def _render_title_slide(presentation: Presentation, slide: TitleSlide) -> None:
        """Renderiza um slide de título (layout 0 — title slide)."""
        layout = presentation.slide_layouts[0]
        s = presentation.slides.add_slide(layout)
        s.shapes.title.text = slide.title
        # Subtítulo opcional: segundo placeholder do layout, se existir.
        if slide.subtitle and len(s.placeholders) > 1:
            s.placeholders[1].text = slide.subtitle

    @staticmethod
    def _render_bullet_slide(presentation: Presentation, slide: BulletSlide) -> None:
        """Renderiza um slide de conteúdo (layout 1 — title and content) com bullets."""
        layout = presentation.slide_layouts[1]
        s = presentation.slides.add_slide(layout)
        s.shapes.title.text = slide.title
        # Placeholder de body: índice 1 quando o layout tem title + content.
        body = None
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break
        if body is None:
            return
        tf = body.text_frame
        # Primeiro bullet já existe no placeholder; reusa para o primeiro item.
        for index, bullet in enumerate(slide.bullets):
            if index == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    @staticmethod
    def _render_image_slide(presentation: Presentation, slide: ImageSlide) -> None:
        """Renderiza um slide com imagem (e título opcional)."""
        layout = presentation.slide_layouts[5]  # title-only
        s = presentation.slides.add_slide(layout)
        if slide.title:
            s.shapes.title.text = slide.title
        _add_image(s, slide.image)

    @staticmethod
    def _render_table_slide(presentation: Presentation, slide: TableSlide) -> None:
        """Renderiza um slide com tabela simples (e título opcional)."""
        layout = presentation.slide_layouts[5]  # title-only
        s = presentation.slides.add_slide(layout)
        if slide.title:
            s.shapes.title.text = slide.title
        _add_table(s, slide.table)


def _add_image(slide, image: ImageRef) -> None:
    """Adiciona `image` ao slide, centralizada, respeitando width_inches."""
    if not Path(image.path).is_file():
        raise RuntimeError(f"Imagem não encontrada: {image.path!r}.")
    left = Inches(1)
    top = Inches(2)
    if image.width_inches is not None:
        slide.shapes.add_picture(
            image.path, left, top, width=Inches(image.width_inches),
        )
    else:
        slide.shapes.add_picture(image.path, left, top)


def _add_table(slide, table: Table) -> None:
    """Adiciona `table` ao slide, no canto superior esquerdo."""
    if not table.rows:
        return
    rows = len(table.rows)
    cols = len(table.rows[0])
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    for row_idx, row in enumerate(table.rows):
        for col_idx, value in enumerate(row):
            cell = shape.table.cell(row_idx, col_idx)
            cell.text = value
            if table.header and row_idx == 0:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
