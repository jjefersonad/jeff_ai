"""Converter semântico HTML slides → PPTX via python-pptx (stdlib HTMLParser).

Reconhece `section.slide` / `div.slide` (classe contendo `slide`). CSS ignorado.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from html.parser import HTMLParser
from io import BytesIO

from pptx import Presentation

from src.application.ports.html_document_converter import HtmlDocumentConverter
from src.domain.shared.errors import DomainError

_SKIP_TAGS = frozenset({"script", "style", "head", "title"})
_VOID_SKIP_TAGS = frozenset({"meta", "link"})
_HEADING_TAGS = frozenset({f"h{i}" for i in range(1, 7)})
_SLIDE_TAGS = frozenset({"section", "div"})


def _has_slide_class(attrs: list[tuple[str, str | None]]) -> bool:
    for key, value in attrs:
        if key.lower() == "class" and value:
            classes = value.split()
            if "slide" in classes:
                return True
    return False


@dataclass
class _SlideDraft:
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    paragraphs: list[str] = field(default_factory=list)


class _HtmlSlideExtractor(HTMLParser):
    """Extrai slides de `section.slide` / `div.slide`."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.slides: list[_SlideDraft] = []
        self._skip_depth = 0
        self._slide_depth = 0
        self._current: _SlideDraft | None = None
        self._capture_tag: str | None = None
        self._parts: list[str] = []
        self._in_list = False

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        lower = tag.lower()
        if lower in _VOID_SKIP_TAGS:
            return
        if lower in _SKIP_TAGS:
            self._skip_depth += 1
            return
        if self._skip_depth:
            return

        if lower in _SLIDE_TAGS and _has_slide_class(attrs) and self._slide_depth == 0:
            self._slide_depth = 1
            self._current = _SlideDraft()
            return
        if self._slide_depth == 0:
            return
        if lower in _SLIDE_TAGS:
            self._slide_depth += 1

        if lower in _HEADING_TAGS and self._current is not None:
            self._flush_capture()
            self._capture_tag = "heading"
            self._parts = []
            return
        if lower == "li" and self._current is not None:
            self._flush_capture()
            self._capture_tag = "li"
            self._parts = []
            return
        if lower in {"ul", "ol"}:
            self._in_list = True
            return
        if lower == "p" and self._current is not None and not self._in_list:
            self._flush_capture()
            self._capture_tag = "p"
            self._parts = []

    def handle_endtag(self, tag: str) -> None:
        lower = tag.lower()
        if lower in _SKIP_TAGS:
            if self._skip_depth:
                self._skip_depth -= 1
            return
        if self._skip_depth:
            return

        if self._slide_depth == 0:
            return

        if lower in _HEADING_TAGS or lower == "li" or lower == "p":
            self._flush_capture()
        if lower in {"ul", "ol"}:
            self._in_list = False

        if lower in _SLIDE_TAGS:
            self._slide_depth -= 1
            if self._slide_depth == 0 and self._current is not None:
                self._flush_capture()
                if self._current.title.strip():
                    self.slides.append(self._current)
                self._current = None

    def handle_data(self, data: str) -> None:
        if self._skip_depth or self._capture_tag is None:
            return
        self._parts.append(data)

    def _flush_capture(self) -> None:
        if self._capture_tag is None or self._current is None:
            self._capture_tag = None
            self._parts = []
            return
        text = "".join(self._parts).strip()
        tag = self._capture_tag
        self._capture_tag = None
        self._parts = []
        if not text:
            return
        if tag == "heading":
            if not self._current.title:
                self._current.title = text
            else:
                self._current.paragraphs.append(text)
        elif tag == "li":
            self._current.bullets.append(text)
        elif tag == "p":
            if not self._current.subtitle and not self._current.bullets:
                self._current.subtitle = text
            else:
                self._current.paragraphs.append(text)


class HtmlPptxConverter(HtmlDocumentConverter):
    """HTML com section/div.slide → bytes PPTX (python-pptx)."""

    async def convert(self, *, html: str, css: str | None, kind: str) -> bytes:
        _ = css
        if kind != "pptx":
            raise DomainError(
                f"HtmlPptxConverter só suporta kind='pptx', recebeu {kind!r}."
            )
        if not html or not html.strip():
            raise DomainError("HTML vazio: nada para converter para PPTX.")

        extractor = _HtmlSlideExtractor()
        extractor.feed(html)
        extractor.close()

        if not extractor.slides:
            raise DomainError(
                "HTML sem slides utilizáveis (section.slide / div.slide com título)."
            )

        presentation = Presentation()
        for draft in extractor.slides:
            _render_slide(presentation, draft)

        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()


def _render_slide(presentation: Presentation, draft: _SlideDraft) -> None:
    bullets = list(draft.bullets)
    bullets.extend(draft.paragraphs)
    if bullets:
        layout = presentation.slide_layouts[1]  # title and content
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = draft.title
        body = slide.placeholders[1].text_frame
        body.clear()
        for idx, item in enumerate(bullets):
            if idx == 0:
                body.paragraphs[0].text = item
            else:
                p = body.add_paragraph()
                p.text = item
                p.level = 0
        return

    layout = presentation.slide_layouts[0]  # title slide
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = draft.title
    if draft.subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = draft.subtitle
